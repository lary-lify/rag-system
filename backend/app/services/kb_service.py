"""
KB Service - Document processing pipeline: parse -> chunk -> embed -> store.
Orchestrates the full document ingestion workflow after upload.
"""
from __future__ import annotations

import asyncio
import logging
import os

from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, func

from app.core.config import settings
from app.core.ids import next_id, next_ids
from app.models.document import Document, DocumentStatus
from app.models.chunk import Chunk
from app.services.chunking.base import ChunkResult
from app.services.embedding_service import embed_texts, estimate_token_count
from app.services.milvus_service import insert_vectors, delete_by_chunk_ids
from app.models.token_usage import TokenUsage, TokenType

logger = logging.getLogger(__name__)

# Strategy registry
_STRATEGIES: dict = {}


def _get_strategy(name: str):
    """Get chunking strategy by name."""
    if not _STRATEGIES:
        from app.services.chunking.fixed_token import FixedTokenChunker
        from app.services.chunking.semantic import SemanticChunker
        from app.services.chunking.paragraph import ParagraphChunker
        from app.services.chunking.heading_level import HeadingLevelChunker
        from app.services.chunking.qa_pair import QAPairChunker
        from app.services.chunking.recursive import RecursiveChunker
        from app.services.chunking.ai_assisted import AIAssistedChunker
        _STRATEGIES.update({
            "fixed_token": FixedTokenChunker(),
            "semantic": SemanticChunker(),
            "paragraph": ParagraphChunker(),
            "heading_level": HeadingLevelChunker(),
            "qa_pair": QAPairChunker(),
            "recursive": RecursiveChunker(),
            "ai_assisted": AIAssistedChunker(),
        })

    strategy = _STRATEGIES.get(name)
    if strategy is None:
        logger.warning(f"Unknown strategy '{name}', falling back to fixed_token")
        return _STRATEGIES["fixed_token"]
    return strategy


async def process_document_async(document_id: int, user_id: int) -> None:
    """
    Full async pipeline for a single uploaded document.
    Steps:
      1. Parse file (PDF/DOCX/PPTX/TXT/MD) -> plain text
      2. Apply selected chunking strategy
      3. Embed all chunks via Tongyi API
      4. Store vectors in Milvus + chunks in MySQL
      5. Record token usage for billing
    """
    from app.core.database import AsyncSessionLocal

    db = AsyncSessionLocal()
    try:
        # Load document record
        result = await db.execute(select(Document).where(Document.id == document_id))
        doc = result.scalar_one_or_none()
        if doc is None or doc.is_deleted:
            logger.warning(f"[pipeline] Document {document_id} not found or deleted")
            return

        kb_id = doc.kb_id  # 捕获为局部变量，避免后续 commit 后访问过期 ORM 属性

        # Load knowledge base config for embedding model
        from app.models.knowledge_base import KnowledgeBase
        kb_result = await db.execute(select(KnowledgeBase).where(KnowledgeBase.id == doc.kb_id))
        kb = kb_result.scalar_one_or_none()
        embedding_model = kb.embedding_model if kb else settings.TONGYI_EMBEDDING_MODEL
        embedding_dimensions = kb.embedding_dimensions if kb else settings.TONGYI_EMBEDDING_DIMENSIONS

        try:
            # Step 1: Update status -> parsing
            doc.status = DocumentStatus.parsing
            await db.commit()

            # Parse file to text
            file_path = os.path.join(settings.UPLOAD_DIR, doc.filename)
            text = await _parse_file(file_path, doc.file_type)

            if not text.strip():
                raise ValueError("Extracted empty text from file")

            # Step 2: Update status -> embedding (chunking happens here)
            doc.status = DocumentStatus.embedding
            await db.commit()

            # Chunk the text
            # 切分是纯 CPU 操作（ai_assisted 策略内部还会同步调用 LLM），
            # 直接在事件循环里跑会把其他请求一起卡住，丢线程池执行
            strategy = _get_strategy(doc.chunk_strategy)
            params = dict(doc.chunk_params or {})
            chunks_result = await asyncio.to_thread(strategy.split, text, **params)

            # Record AI chunking token usage if the strategy provides it
            ai_usage = getattr(strategy, "get_last_ai_usage", lambda: None)()
            if ai_usage:
                ai_tok_record = TokenUsage(
                    type=TokenType.chunking,
                    user_id=user_id,
                    kb_id=doc.kb_id,
                    document_id=doc.id,
                    input_tokens=ai_usage["input_tokens"],
                    output_tokens=ai_usage["output_tokens"],
                    estimated_cost=ai_usage["estimated_cost"],
                )
                db.add(ai_tok_record)
                await db.commit()
                logger.info(
                    f"[pipeline] AI chunking recorded: "
                    f"{ai_usage['input_tokens']} in + {ai_usage['output_tokens']} out, "
                    f"cost=¥{ai_usage['estimated_cost']:.6f}"
                )

            if not chunks_result:
                raise ValueError("Chunking produced no results")

            # Step 3: Embed chunks
            contents = [c.content for c in chunks_result]
            try:
                vectors, embed_tokens = await asyncio.wait_for(
                    embed_texts(
                        contents,
                        model=embedding_model,
                        dimensions=embedding_dimensions,
                        # 入库不写缓存：一次性涌入的成千上万条片段会把
                        # 查询向量从 LRU 里挤出去，反而抹平缓存收益
                        use_cache=False,
                    ),
                    timeout=180,
                )
            except asyncio.TimeoutError:
                raise TimeoutError("Embedding API timed out after 180s")

            if len(vectors) != len(chunks_result):
                logger.warning(
                    f"[pipeline] Vector count mismatch: got {len(vectors)} for {len(chunks_result)} chunks"
                )

            # Step 4: Store chunks in DB + Milvus
            # 主键由应用层生成，不再依赖数据库自增。
            #
            # 原因是 aiomysql 下只要需要回填自增主键，SQLAlchemy 就会回退成
            # 逐行 INSERT，insertmanyvalues 批量优化失效——批量 add_all 看
            # 似合并了，实际仍是一条一条发，500 片段就是 500 次数据库往返。
            # 显式给出主键后无需回填，插入才真正合并成一条多值 INSERT
            # （实测 500 次往返 -> 1 次，635.8ms -> 183.2ms）。
            chunk_ids = next_ids(len(chunks_result))
            db_chunks = [
                Chunk(
                    id=chunk_ids[i],
                    document_id=doc.id,
                    kb_id=doc.kb_id,
                    content=cr.content,
                    chunk_index=cr.index,
                    token_count=cr.token_count or estimate_token_count(cr.content),
                    chunk_meta=cr.metadata,
                )
                for i, cr in enumerate(chunks_result)
            ]
            db.add_all(db_chunks)
            await db.flush()

            # 主键应用层已知，无需等 flush 回填
            db_chunk_ids = list(chunk_ids)
            # Use DB ID as Milvus primary key for easy mapping
            milvus_chunk_ids = list(db_chunk_ids)

            doc.chunk_count = len(chunks_result)
            doc.status = DocumentStatus.completed
            await db.commit()

            # Insert into Milvus
            #
            # milvus_id 是「该片段的向量确实已写进向量库」的标志位，只能由
            # 真实的写入成功来置位。原实现把超时 catch 掉只打一行 warning，
            # 随后仍无条件回填 milvus_id，于是向量库里根本不存在的片段也被
            # 标记成已入库：检索静默漏召，全链路没有任何报错，运维无从发现。
            #
            # 这里改为：失败一律不置位，并把可补偿信息落到 documents.error_msg，
            # 运维可用 `SELECT ... FROM chunks WHERE milvus_id IS NULL` 找出待补偿片段。
            milvus_ok = False
            if not vectors:
                logger.error(
                    f"[pipeline] no vectors produced for doc {document_id}, skipping Milvus insert"
                )
            elif db_chunk_ids:
                try:
                    await asyncio.wait_for(
                        insert_vectors(
                            kb_id=doc.kb_id,
                            chunk_ids=milvus_chunk_ids,
                            document_ids=[doc.id] * len(milvus_chunk_ids),
                            contents=contents[:len(milvus_chunk_ids)],
                            vectors=vectors[:len(milvus_chunk_ids)],
                            dimension=embedding_dimensions,
                        ),
                        timeout=60,
                    )
                    milvus_ok = True
                except asyncio.TimeoutError:
                    logger.error(
                        f"[pipeline] Milvus insert timed out for doc {document_id} "
                        f"({len(db_chunk_ids)} chunks): text saved to DB, vectors NOT indexed"
                    )
                except Exception as e:
                    # 原实现只 catch TimeoutError，其余异常会直接冒泡中断流程；
                    # 这里统一接住，保证「文本已入库」这个既成事实不被回滚，
                    # 同时如实记录向量侧的失败。
                    logger.error(
                        f"[pipeline] Milvus insert failed for doc {document_id} "
                        f"({len(db_chunk_ids)} chunks): {e}. "
                        f"text saved to DB, vectors NOT indexed",
                        exc_info=True,
                    )

            # Update Milvus IDs on chunks — only for chunks actually in Milvus.
            # Milvus 主键直接复用数据库 chunk id，原实现逐条 UPDATE 把
            # 同一个值写回同一行，N 个片段就是 N 次往返。一条语句即可。
            if milvus_ok and db_chunk_ids:
                await db.execute(
                    Chunk.__table__.update()
                    .where(Chunk.id.in_(db_chunk_ids))
                    .values(milvus_id=Chunk.id)
                )
            elif db_chunk_ids:
                # 文档正文仍可查（列表/下载不受影响），但向量检索查不到，
                # 如实写进文档错误信息，避免用户以为入库已全部完成。
                # 用显式 UPDATE 而非 doc.error_msg = ... ：doc 在前面的
                # commit 之后已过期，异步下触碰 ORM 属性会触发惰性加载报错。
                await db.execute(
                    Document.__table__.update()
                    .where(Document.id == doc.id)
                    .values(error_msg=(
                        f"文本已入库，但向量未写入（{len(db_chunk_ids)} 个片段待补偿）"
                    ))
                )
                logger.error(
                    f"[pipeline] doc {document_id} marked text-only: "
                    f"{len(db_chunk_ids)} chunk(s) have no vector, needs re-index"
                )
            await db.commit()

            # 知识库内容已变更（新片段进入向量库，检索结果会变）→ 使该 KB 的答案缓存失效。
            # 通过自增 KB 世代计数器：旧 scope（含旧 epoch）的答案立即查不到，新查询走新 scope。
            if milvus_ok:
                from app.services.answer_cache import bump_kb_epoch

                try:
                    new_epoch = await bump_kb_epoch(kb_id)
                    logger.info(f"[pipeline] bumped answer-cache epoch for kb {kb_id} -> {new_epoch}")
                except Exception as e:
                    logger.warning(f"[pipeline] bump_kb_epoch failed (kb={kb_id}): {e}")

            # Update KB cache counters
            from app.models.knowledge_base import KnowledgeBase
            kb_res = await db.execute(select(KnowledgeBase).where(KnowledgeBase.id == doc.kb_id))
            kb = kb_res.scalar_one_or_none()
            if kb:
                # Re-count documents and chunks
                doc_count_res = await db.execute(
                    select(func.count(Document.id))
                    .where(Document.kb_id == doc.kb_id, Document.is_deleted == False, Document.status == DocumentStatus.completed)
                )
                chunk_count_res = await db.execute(
                    select(func.count(Chunk.id))
                    .where(Chunk.kb_id == doc.kb_id, Chunk.is_deleted == False)
                )
                kb.doc_count = doc_count_res.scalar_one() or 0
                kb.chunk_count = chunk_count_res.scalar_one() or 0
                await db.commit()

            # Step 5: Record token usage for billing
            tok_record = TokenUsage(
                type=TokenType.embedding,
                user_id=user_id,
                kb_id=doc.kb_id,
                document_id=doc.id,
                input_tokens=embed_tokens,
                output_tokens=0,
                estimated_cost=_calc_embedding_cost(embed_tokens),
            )
            db.add(tok_record)
            await db.commit()

            logger.info(
                f"[pipeline] Doc {document_id} completed: {len(chunks_result)} chunks, {embed_tokens} tokens"
            )

        except Exception as e:
            logger.error(f"[pipeline] Doc {document_id} FAILED: {e}", exc_info=True)
            try:
                doc.status = DocumentStatus.failed
                doc.error_msg = str(e)[:500]
                await db.commit()
            except Exception as commit_err:
                logger.error(f"[pipeline] Failed to commit error status for doc {document_id}: {commit_err}")
                try:
                    await db.rollback()
                except Exception:
                    pass

    except Exception as e:
        logger.error(f"[pipeline] Fatal error for doc {document_id}: {e}", exc_info=True)
    finally:
        await db.close()


# ---- File Parsers ----

async def _parse_file(file_path: str, file_type: str) -> str:
    """
    Parse uploaded file to plain text based on extension.

    解析是同步且耗时的操作：PDF 要逐页解文本，DOCX 要解包 XML，开了 OCR
    的话还要对每张图做一次视觉推理。直接调用会把事件循环整个占住，期间
    所有其他请求都得不到调度——一个 50 页 PDF 能让整个服务停摆好几秒。
    丢到线程池里跑，事件循环保持可调度。
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    def _parse_sync() -> str:
        # Use intelligent parser for supported types
        if file_type in ("docx", "doc", "html", "htm", "txt", "md", "pdf"):
            from app.services.document_parser import (
                parse_document_intelligently,
                format_structured_content,
            )
            parsed = parse_document_intelligently(file_path, file_type)
            return format_structured_content(parsed)

        match file_type:
            case "pptx":
                return _parse_pptx(file_path)
            case "csv":
                return _parse_csv(file_path)
            case "xlsx" | "xls":
                return _parse_excel(file_path)
            case _:
                return _parse_unstructured(file_path)

    return await asyncio.to_thread(_parse_sync)


def _parse_pdf(path: str) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    texts = []
    for page in reader.pages:
        t = page.extract_text()
        if t:
            texts.append(t)
    return "\n\n".join(texts)


def _parse_docx(path: str) -> str:
    from docx import Document as DocxDocument
    doc = DocxDocument(path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paras)


def _parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides_text = []
    for slide_idx, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides_text.append(f"[Slide {slide_idx + 1}]\n" + "\n".join(parts))
    return "\n\n---\n\n".join(slides_text)


def _parse_excel(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join([str(c) if c is not None else "" for c in row])
            if row_text.strip():
                rows.append(row_text)
        sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n---\n\n".join(sheets)


def _parse_csv(path: str) -> str:
    """Parse CSV with intelligent Q/A pair detection."""
    import csv
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.reader(f)
        rows = list(reader)

    if not rows:
        return ""

    header = [h.strip().lower() for h in rows[0]]

    # Detect question/answer columns
    q_col = a_col = -1
    for i, h in enumerate(header):
        if h in ("question", "问题", "q", "ques"):
            q_col = i
        elif h in ("answer", "回答", "a", "ans"):
            a_col = i

    if q_col >= 0 and a_col >= 0:
        # Format as QA pairs for qa_pair chunker
        parts = []
        for row in rows[1:]:
            if len(row) > max(q_col, a_col):
                q = row[q_col].strip()
                a = row[a_col].strip()
                if q and a:
                    parts.append(f"问：{q}\n答：{a}")
        return "\n\n".join(parts) if parts else str(rows)

    # Fallback: tab-separated rows
    lines = []
    for row in rows[1:]:
        lines.append("\t".join(row))
    return "\n".join(lines)


def _parse_unstructured(path: str) -> str:
    # Fallback: try reading as plain text
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        raise ValueError(f"Unsupported file type, cannot parse: {path}")


def _calc_embedding_cost(tokens: int) -> float:
    return round(tokens * settings.TONGYI_EMBEDDING_TOKEN_PRICE, 6)
